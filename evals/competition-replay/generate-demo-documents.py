#!/usr/bin/env python3
"""生成可反复上传的 Word、PPT、Excel 和 PDF 比赛演示附件。"""

from __future__ import annotations

import subprocess
from datetime import datetime, timezone
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "demo-materials"
FIXED_TIME = datetime(2026, 8, 12, tzinfo=timezone.utc)


def generate_word() -> Path:
    """生成包含主体、财务和风险三类主张的 Word 材料，并固定核心属性时间。"""
    target = OUTPUT / "06-模拟企业尽调材料.docx"
    document = Document()
    document.core_properties.title = "云岚数据（苏州）有限公司尽调材料"
    document.core_properties.created = FIXED_TIME
    document.core_properties.modified = FIXED_TIME
    document.add_heading("云岚数据（苏州）有限公司尽调摘要", level=1)
    document.add_paragraph("纯模拟材料：企业名称、编码和全部数据均为虚构，仅用于内部比赛演示。")
    document.add_heading("待核验事实", level=2)
    document.add_paragraph("1. 统一社会信用代码为 91320500MA2DEMO006。")
    document.add_paragraph("2. 2025 年营业收入为 9800 万元。")
    document.add_paragraph("3. 2025 年不存在任何行政处罚或失信记录。")
    document.save(target)
    return target


def generate_powerpoint() -> Path:
    """生成两页 PPT，供浏览器路径验证 slide locator 与多页文本提取。"""
    target = OUTPUT / "07-模拟企业融资说明.pptx"
    presentation = Presentation()
    presentation.core_properties.title = "澄海智能装备有限公司融资说明"
    presentation.core_properties.created = FIXED_TIME
    presentation.core_properties.modified = FIXED_TIME
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = "澄海智能装备有限公司"
    title_slide.placeholders[1].text = "纯模拟内部比赛材料，不对应任何真实企业"
    facts_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    facts_slide.shapes.title.text = "经营事实"
    facts_slide.placeholders[1].text = (
        "统一社会信用代码：91330200MA2DEMO007\n"
        "2025 年营业收入：1.25 亿元\n"
        "截至 2025 年末拥有有效专利 18 件"
    )
    presentation.save(target)
    return target


def generate_excel() -> Path:
    """生成带公式的 Excel 台账，验证单元格定位、单位和公式文本都被保留。"""
    target = OUTPUT / "08-模拟企业财务台账.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "经营指标"
    sheet.append(["纯模拟数据", "以下企业和指标不对应任何真实主体"])
    sheet.append(["企业名称", "北辰工业软件有限公司"])
    sheet.append(["统一社会信用代码", "91310100MA2DEMO008"])
    sheet.append(["年份", "营业收入万元", "研发投入万元", "研发投入强度"])
    sheet.append([2024, 7600, 760, "=C5/B5"])
    sheet.append([2025, 10200, 1326, "=C6/B6"])
    sheet["D5"].number_format = "0.0%"
    sheet["D6"].number_format = "0.0%"
    workbook.properties.created = FIXED_TIME.replace(tzinfo=None)
    workbook.properties.modified = FIXED_TIME.replace(tzinfo=None)
    workbook.save(target)
    return target


def generate_pdf(word_path: Path) -> Path:
    """复用同一 Word 内容经 LibreOffice 生成 PDF，避免维护两份含义不同的模拟材料。"""
    subprocess.run(
        [
            "libreoffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(OUTPUT),
            str(word_path),
        ],
        check=True,
        capture_output=True,
        text=True,
    )
    return word_path.with_suffix(".pdf")


def main() -> None:
    """一次性重建四份二进制附件，并拒绝静默漏产文件。"""
    OUTPUT.mkdir(parents=True, exist_ok=True)
    word_path = generate_word()
    generated = [word_path, generate_powerpoint(), generate_excel(), generate_pdf(word_path)]
    missing = [str(path) for path in generated if not path.is_file() or path.stat().st_size == 0]
    if missing:
        raise RuntimeError(f"演示附件生成失败：{missing}")
    for path in generated:
        print(f"GENERATED {path.name} {path.stat().st_size} bytes")


if __name__ == "__main__":
    main()
