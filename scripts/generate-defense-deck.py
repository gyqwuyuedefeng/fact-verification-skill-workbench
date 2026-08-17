#!/usr/bin/env python3
"""生成“火石验真官”竞赛答辩 PPT。"""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "submission/答辩材料/火石验真官-企业材料事实核验-答辩演示.pptx"
SCREENSHOT_DIR = ROOT / "submission/评价结果/演示截图"

SLIDE_W = 13.333
SLIDE_H = 7.5

NAVY = "061522"
NAVY_2 = "0A1D2D"
CARD = "0C2234"
CARD_2 = "102B40"
TEAL = "4FD1C5"
TEAL_DARK = "0D7D79"
CYAN = "72E8DF"
WHITE = "F6FAFC"
MUTED = "9FB6C9"
MUTED_2 = "6F8AA0"
YELLOW = "F6C85F"
RED = "FF6B72"
GREEN = "65D6A8"
FONT = "Microsoft YaHei"
MONO = "Consolas"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def add_rect(slide, x, y, w, h, fill, line=None, radius=True, transparency=0):
    """添加统一圆角卡片；所有坐标使用英寸，便于保持 16:9 栅格。"""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = transparency
    shape.line.color.rgb = rgb(line or fill)
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    size=20,
    color=WHITE,
    bold=False,
    font=FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
    margin=0.02,
):
    """添加单一风格文本框，主动关闭自动换号以保持版式可控。"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_rule(slide, x, y, w, color=TEAL_DARK):
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.018)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(color)
    line.line.fill.background()


def add_title(slide, eyebrow: str, title: str, subtitle: str | None = None):
    add_text(slide, eyebrow.upper(), 0.72, 0.34, 4.0, 0.30, size=10, color=TEAL, bold=True, font=MONO)
    add_text(slide, title, 0.72, 0.72, 11.9, 0.55, size=27, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.74, 1.29, 11.7, 0.34, size=12, color=MUTED)


def add_footer(slide, number: int):
    add_rule(slide, 0.72, 7.08, 11.89, color="18384C")
    add_text(
        slide,
        "火石验真官 · AI即岗位，结果即作品",
        0.74,
        7.12,
        5.2,
        0.20,
        size=8,
        color=MUTED_2,
        font=MONO,
    )
    add_text(
        slide,
        f"{number:02d} / 07",
        11.4,
        7.12,
        1.18,
        0.20,
        size=8,
        color=TEAL,
        bold=True,
        font=MONO,
        align=PP_ALIGN.RIGHT,
    )


def add_chip(slide, text, x, y, w, *, fill=CARD_2, color=CYAN, line=TEAL_DARK):
    add_rect(slide, x, y, w, 0.36, fill, line)
    add_text(slide, text, x + 0.04, y + 0.01, w - 0.08, 0.31, size=10, color=color, bold=True, align=PP_ALIGN.CENTER)


def add_card(slide, number, title, body, x, y, w, h, *, accent=TEAL):
    add_rect(slide, x, y, w, h, CARD, "173A50")
    add_text(slide, number, x + 0.22, y + 0.20, 0.44, 0.28, size=10, color=accent, bold=True, font=MONO)
    add_text(slide, title, x + 0.22, y + 0.58, w - 0.44, 0.40, size=19, bold=True)
    add_text(slide, body, x + 0.22, y + 1.07, w - 0.44, h - 1.30, size=12, color=MUTED, valign=MSO_ANCHOR.TOP)


def add_arrow(slide, x, y, w=0.50):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(0.55)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(TEAL_DARK)
    shape.line.fill.background()


def add_picture_contain(slide, path: Path, x, y, w, h):
    """按原比例完整放入截图，避免拉伸业务界面。"""
    with Image.open(path) as image:
        image_ratio = image.width / image.height
    box_ratio = w / h
    if image_ratio >= box_ratio:
        width = w
        height = w / image_ratio
        left = x
        top = y + (h - height) / 2
    else:
        height = h
        width = h * image_ratio
        left = x + (w - width) / 2
        top = y
    return slide.shapes.add_picture(
        str(path), Inches(left), Inches(top), Inches(width), Inches(height)
    )


def add_picture_crop(slide, path: Path, x, y, w, h, *, focus_y=0.5):
    """裁切超长页面截图，只调整 PPT 裁切参数，不改动原始证据图片。"""
    with Image.open(path) as image:
        image_ratio = image.width / image.height
    box_ratio = w / h
    picture = slide.shapes.add_picture(
        str(path), Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if image_ratio > box_ratio:
        visible = box_ratio / image_ratio
        crop = 1 - visible
        picture.crop_left = crop / 2
        picture.crop_right = crop / 2
    else:
        visible = image_ratio / box_ratio
        crop = 1 - visible
        picture.crop_top = crop * focus_y
        picture.crop_bottom = crop * (1 - focus_y)
    return picture


def new_slide(presentation):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = rgb(NAVY)
    return slide


def slide_cover(presentation):
    slide = new_slide(presentation)
    add_rect(slide, 0.0, 0.0, 13.333, 7.5, NAVY, NAVY, radius=False)
    add_rect(slide, 8.62, -0.55, 5.8, 5.8, TEAL_DARK, TEAL_DARK, transparency=78)
    add_text(slide, "818 AI 全员竞赛 · 模型与 Agent", 0.78, 0.55, 5.4, 0.35, size=11, color=TEAL, bold=True, font=MONO)
    add_text(slide, "企业材料事实核验", 0.78, 1.34, 11.2, 0.90, size=39, bold=True)
    add_text(slide, "让 Skill 能被验证、被发布，也能被安全回滚", 0.82, 2.39, 10.9, 0.58, size=21, color=MUTED)
    for text, x, w in (("AGENT", 0.82, 1.26), ("SKILL", 2.22, 1.20), ("MCP", 3.56, 1.08), ("EVAL", 4.78, 1.10)):
        add_chip(slide, text, x, 3.35, w)
    add_rule(slide, 0.82, 4.25, 5.15)
    add_text(slide, "队名", 0.82, 4.58, 0.78, 0.32, size=10, color=MUTED_2, font=MONO)
    add_text(slide, "火石验真官", 1.62, 4.53, 3.20, 0.38, size=18, color=CYAN, bold=True)
    add_text(slide, "成员", 0.82, 5.06, 0.78, 0.32, size=10, color=MUTED_2, font=MONO)
    add_text(slide, "高永强 · 高翔宇 · 钱旭", 1.62, 5.02, 5.20, 0.38, size=16, bold=True)
    add_text(slide, "AI即岗位，结果即作品", 8.73, 5.80, 3.65, 0.40, size=14, color=CYAN, bold=True, align=PP_ALIGN.RIGHT)
    add_text(slide, "可运行 · 可验证 · 可复制", 8.73, 6.24, 3.65, 0.28, size=10, color=MUTED, font=MONO, align=PP_ALIGN.RIGHT)
    add_footer(slide, 1)


def slide_problem(presentation):
    slide = new_slide(presentation)
    add_title(slide, "01 / BUSINESS PROBLEM", "为什么做", "企业材料进入报告、招商或尽调前，必须先回答“这句话是真的吗？”")
    add_card(slide, "01", "人工核验慢", "一份材料包含多个主体、指标与风险事实，逐条查询、比对和留痕耗时。", 0.74, 1.88, 3.80, 3.55)
    add_card(slide, "02", "通用模型会猜", "模型能读懂文字，但没有企业证据时，容易把语言流畅误当成事实正确。", 4.76, 1.88, 3.80, 3.55, accent=YELLOW)
    add_card(slide, "03", "结论不可追溯", "只有“对或错”还不够，业务需要知道原文位置、证据来源与人工边界。", 8.78, 1.88, 3.80, 3.55, accent=RED)
    add_rect(slide, 0.74, 5.73, 11.84, 0.78, NAVY_2, TEAL_DARK)
    add_text(slide, "我们要回答：材料中的事实，是 已验证、存在冲突，还是证据不足？", 0.98, 5.84, 11.34, 0.48, size=17, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 2)


def slide_solution(presentation):
    slide = new_slide(presentation)
    add_title(slide, "02 / SOLUTION", "我们做了什么", "不是让模型自由发挥，而是让 Agent 按固定 Skill 调用只读企业证据。")
    x_values = (0.74, 3.83, 6.92, 10.01)
    titles = ("材料输入", "Agent拆分主张", "只读MCP查证", "三类可追溯结论")
    bodies = ("文本或文件", "主体 · 指标 · 时间", "六类企业数据", "验证 · 冲突 · 证据不足")
    accents = (TEAL, CYAN, YELLOW, GREEN)
    for index, (x, title, body, accent) in enumerate(zip(x_values, titles, bodies, accents), start=1):
        add_rect(slide, x, 2.05, 2.34, 2.30, CARD, "173A50")
        add_text(slide, f"0{index}", x + 0.19, 2.25, 0.48, 0.28, size=10, color=accent, bold=True, font=MONO)
        add_text(slide, title, x + 0.19, 2.74, 1.96, 0.46, size=17, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.19, 3.43, 1.96, 0.42, size=11, color=MUTED, align=PP_ALIGN.CENTER)
        if index < 4:
            add_arrow(slide, x + 2.48, 2.91)
    add_text(slide, "可信结果必须能解释", 0.74, 4.84, 3.00, 0.34, size=12, color=MUTED, bold=True)
    for text, x, w in (("原文定位", 0.74, 2.13), ("证据 recordId", 3.06, 2.35), ("人工介入边界", 5.60, 2.43)):
        add_chip(slide, text, x, 5.39, w)
    add_rect(slide, 8.50, 4.77, 4.08, 1.17, NAVY_2, TEAL_DARK)
    add_text(slide, "一个 Agent + 一个 Skill\n+ 六个只读 MCP 工具", 8.76, 4.91, 3.56, 0.82, size=16, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 3)


def slide_demo(presentation):
    slide = new_slide(presentation)
    add_title(slide, "03 / LIVE DEMO", "现场演示：一条真实材料如何被核验")
    for text, x, w in (("01 输入快照", 0.74, 1.55), ("02 执行轨迹", 2.45, 1.55), ("03 核验主张", 4.16, 1.55)):
        add_chip(slide, text, x, 1.38, w)
    add_rect(slide, 0.72, 1.91, 11.90, 4.47, CARD, "1C455B")
    add_picture_contain(
        slide,
        SCREENSHOT_DIR / "01-普通用户事实核验-解析轨迹与主张.png",
        0.88,
        2.06,
        11.58,
        4.04,
    )
    add_rect(slide, 1.52, 6.20, 10.29, 0.48, NAVY_2, TEAL_DARK)
    add_text(slide, "材料原文、工具调用和结论证据同时留痕，结果可以复核。", 1.68, 6.25, 9.98, 0.30, size=12, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 4)


def slide_evaluation(presentation):
    slide = new_slide(presentation)
    add_title(slide, "04 / REPRODUCIBLE EVALUATION", "同条件评测：只改变 Skill", "同一数据集、同一输入、同一模型参数、同一工具与证据快照。")
    add_rect(slide, 0.74, 1.82, 6.08, 4.59, CARD, "173A50")
    add_text(slide, "准确样本数", 1.05, 2.13, 1.60, 0.30, size=11, color=MUTED, font=MONO)
    add_text(slide, "7/30", 1.02, 2.62, 1.48, 0.70, size=30, color=RED, bold=True, font=MONO, align=PP_ALIGN.CENTER)
    add_arrow(slide, 2.56, 2.69, 0.52)
    add_text(slide, "27/30", 3.12, 2.62, 1.68, 0.70, size=30, color=YELLOW, bold=True, font=MONO, align=PP_ALIGN.CENTER)
    add_arrow(slide, 4.91, 2.69, 0.52)
    add_text(slide, "30/30", 5.40, 2.62, 1.18, 0.70, size=27, color=GREEN, bold=True, font=MONO, align=PP_ALIGN.CENTER)
    add_text(slide, "通用模型\nBASELINE", 1.02, 3.41, 1.48, 0.62, size=10, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, "初始 Skill\nSTABLE", 3.15, 3.41, 1.68, 0.62, size=10, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, "优化 Skill\nCANDIDATE", 5.31, 3.41, 1.38, 0.62, size=10, color=MUTED, align=PP_ALIGN.CENTER)
    add_rule(slide, 1.03, 4.32, 5.50, color="21475D")
    add_text(slide, "八项门禁同时判断", 1.03, 4.59, 2.10, 0.32, size=12, color=CYAN, bold=True)
    add_text(slide, "准确率 · 完成率 · 稳定性 · 人工介入率\n执行成功 · 样本齐全 · 同条件 · 结果可追溯", 1.03, 5.05, 5.50, 0.82, size=13, color=WHITE, bold=True)
    add_rect(slide, 7.09, 1.82, 5.49, 4.59, CARD, "173A50")
    add_picture_crop(
        slide,
        SCREENSHOT_DIR / "02-三版本同条件评测与门禁.png",
        7.29,
        2.03,
        5.09,
        4.17,
        focus_y=0.04,
    )
    add_footer(slide, 5)


def slide_release(presentation):
    slide = new_slide(presentation)
    add_title(slide, "05 / SAFE RELEASE", "安全发布：用户始终只看到 Stable", "Candidate 在后台影子运行，同一请求双跑，但候选结果不返回给普通用户。")
    add_rect(slide, 0.74, 1.83, 5.48, 4.67, CARD, "173A50")
    steps = (
        ("01", "Candidate 注册", "冻结内容与版本 hash"),
        ("02", "影子灰度", "真实请求双跑、不影响用户"),
        ("03", "人工 PASS", "查看差异与失败原因"),
        ("04", "晋升 / 回滚", "切换 Stable 或恢复上一版"),
    )
    for index, (number, title, body) in enumerate(steps):
        y = 2.14 + index * 0.96
        add_rect(slide, 1.02, y, 0.53, 0.53, TEAL_DARK, TEAL_DARK)
        add_text(slide, number, 1.05, y + 0.02, 0.47, 0.46, size=10, color=WHITE, bold=True, font=MONO, align=PP_ALIGN.CENTER)
        add_text(slide, title, 1.77, y - 0.02, 1.90, 0.33, size=14, bold=True)
        add_text(slide, body, 1.77, y + 0.34, 3.92, 0.28, size=10, color=MUTED)
        if index < 3:
            add_rule(slide, 1.27, y + 0.55, 0.02, color=TEAL_DARK)
    for text, x, color in (("一致 1", 1.02, GREEN), ("差异 0", 2.63, CYAN), ("可回滚", 4.24, YELLOW)):
        add_chip(slide, text, x, 5.86, 1.42, color=color)
    add_rect(slide, 6.50, 1.83, 6.08, 4.67, CARD, "173A50")
    add_picture_contain(
        slide,
        SCREENSHOT_DIR / "03-影子晋升与回滚历史.png",
        6.69,
        2.00,
        5.70,
        4.32,
    )
    add_footer(slide, 6)


def slide_plan(presentation):
    slide = new_slide(presentation)
    add_title(slide, "06 / NEXT 30 DAYS", "赛后30天：沉淀为火石智能体通用能力", "目标不是只交付一个 Skill，而是让任意 Skill 都具备版本、评测、灰度与回滚能力。")
    stages = (
        ("1—7天", "统一治理合同", "版本状态、评测口径、发布边界"),
        ("8—14天", "版本与评测", "接入现有智能体管理端"),
        ("15—21天", "通用影子灰度", "真实流量观察与一键回滚"),
        ("22—30天", "第二 Skill 验收", "验证能力不依赖单一场景"),
    )
    colors = (TEAL, CYAN, YELLOW, GREEN)
    for index, ((days, title, body), accent) in enumerate(zip(stages, colors)):
        x = 0.74 + index * 3.02
        add_rect(slide, x, 2.05, 2.78, 2.75, CARD, "173A50")
        add_text(slide, days, x + 0.22, 2.29, 2.32, 0.34, size=12, color=accent, bold=True, font=MONO)
        add_text(slide, title, x + 0.22, 2.91, 2.32, 0.46, size=17, bold=True)
        add_text(slide, body, x + 0.22, 3.56, 2.32, 0.77, size=11, color=MUTED, valign=MSO_ANCHOR.TOP)
    add_rect(slide, 1.53, 5.30, 10.27, 0.88, NAVY_2, TEAL_DARK)
    add_text(slide, "不是只做一个 Skill，而是补齐火石智能体的 Skill 治理闭环", 1.76, 5.46, 9.80, 0.46, size=19, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "谢谢", 5.35, 6.38, 2.64, 0.40, size=21, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 7)


def build() -> None:
    for screenshot in (
        "01-普通用户事实核验-解析轨迹与主张.png",
        "02-三版本同条件评测与门禁.png",
        "03-影子晋升与回滚历史.png",
    ):
        path = SCREENSHOT_DIR / screenshot
        if not path.is_file():
            raise FileNotFoundError(f"缺少演示截图：{path}")

    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_W)
    presentation.slide_height = Inches(SLIDE_H)
    presentation.core_properties.title = "火石验真官——企业材料事实核验"
    presentation.core_properties.subject = "818 AI 全员竞赛答辩演示"
    presentation.core_properties.author = "高永强、高翔宇、钱旭"
    presentation.core_properties.comments = "由项目内生成脚本产生，支持复现。"

    slide_cover(presentation)
    slide_problem(presentation)
    slide_solution(presentation)
    slide_demo(presentation)
    slide_evaluation(presentation)
    slide_release(presentation)
    slide_plan(presentation)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(OUTPUT)
    print(f"DEFENSE_DECK_GENERATED {OUTPUT}")


if __name__ == "__main__":
    build()
