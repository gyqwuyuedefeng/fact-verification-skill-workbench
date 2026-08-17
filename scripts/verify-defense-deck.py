#!/usr/bin/env python3
"""校验竞赛答辩 PPT 的结构合同，避免提交缺页、缺图或泄露内网信息。"""

from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


DEFAULT_DECK = Path("submission/答辩材料/火石验真官-企业材料事实核验-答辩演示.pptx")
REQUIRED_TEXT = (
    "企业材料事实核验",
    "7/30",
    "27/30",
    "30/30",
    "影子灰度",
    "30天",
)
FORBIDDEN_PATTERNS = (
    re.compile(r"192\.168\.\d+\.\d+"),
    re.compile(r"jdbc:", re.IGNORECASE),
    re.compile(r"Bearer\s+\S+", re.IGNORECASE),
    re.compile(r"password\s*[:=]", re.IGNORECASE),
)


def shape_text(shape: object) -> str:
    """提取普通文本框与组合形状中的文字，供提交前合同检查使用。"""
    if getattr(shape, "has_text_frame", False):
        return getattr(shape, "text", "") or ""
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        return "\n".join(shape_text(child) for child in shape.shapes)
    return ""


def verify(deck_path: Path) -> tuple[int, int]:
    """验证页数、比例、关键叙事、截图数量及敏感信息边界。"""
    if not deck_path.is_file():
        raise AssertionError(f"PPT 不存在：{deck_path}")

    presentation = Presentation(deck_path)
    if len(presentation.slides) != 7:
        raise AssertionError(f"PPT 必须为 7 页，实际为 {len(presentation.slides)} 页")

    ratio = presentation.slide_width / presentation.slide_height
    if abs(ratio - 16 / 9) > 0.01:
        raise AssertionError(f"PPT 必须为 16:9，实际宽高比为 {ratio:.3f}")

    all_text: list[str] = []
    picture_count = 0
    for index, slide in enumerate(presentation.slides, start=1):
        slide_text = "\n".join(shape_text(shape) for shape in slide.shapes)
        all_text.append(slide_text)
        if not slide_text.strip():
            raise AssertionError(f"第 {index} 页没有可读文字")
        picture_count += sum(
            1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        )

    merged_text = "\n".join(all_text)
    missing = [text for text in REQUIRED_TEXT if text not in merged_text]
    if missing:
        raise AssertionError(f"PPT 缺少关键叙事：{', '.join(missing)}")
    if picture_count < 3:
        raise AssertionError(f"PPT 至少应包含 3 张真实系统截图，实际为 {picture_count} 张")

    for pattern in FORBIDDEN_PATTERNS:
        if pattern.search(merged_text):
            raise AssertionError(f"PPT 包含不应提交的敏感内容：{pattern.pattern}")

    return len(presentation.slides), picture_count


def main() -> int:
    parser = argparse.ArgumentParser(description="校验竞赛答辩 PPT")
    parser.add_argument("deck", nargs="?", type=Path, default=DEFAULT_DECK)
    args = parser.parse_args()
    try:
        slides, pictures = verify(args.deck)
    except (AssertionError, OSError, ValueError) as error:
        print(f"DEFENSE_DECK_FAIL {error}", file=sys.stderr)
        return 1
    print(f"DEFENSE_DECK_PASS slides={slides} pictures={pictures} ratio=16:9")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
